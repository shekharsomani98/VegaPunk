from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from mistralai import Mistral
import json
import re
import os
import hashlib
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
from pydantic_settings import BaseSettings
from functools import lru_cache
import uvicorn
from pathlib import Path
from pptx import Presentation
from PIL import Image
import base64
import io
import concurrent.futures
from config import config
from utils import load_json, save_json
import matplotlib
matplotlib.use('Agg')  # Use the 'Agg' backend which doesn't require a GUI
import matplotlib.pyplot as plt
from sympy import preview
import shutil


# Create FastAPI app
app = FastAPI(title="Research Paper Prerequisites API")


# Configuration settings
class Settings(BaseSettings):
    MISTRAL_API_KEY: str
    MODEL_NAME: str = "mistral-small-latest"
    EXECUTION_AGENT_ID: Optional[str] = None
    MODEL_NAME_OCR: str = "mistral-ocr-latest"
    
    class Config:
        env_file = ".env"
        extra = 'ignore'

@lru_cache
def get_settings():
    return Settings()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create necessary directories
os.makedirs("data/upload", exist_ok=True)
os.makedirs("data/figures", exist_ok=True)
os.makedirs("data/formulas", exist_ok=True)
os.makedirs("data/metadata", exist_ok=True)
# os.makedirs("data/metadata")

# Define models
class AnalysisResponse(BaseModel):
    prerequisites: dict

class Formula(BaseModel):
    formula: str
    name: str

class Slides(BaseModel):
    title: str
    subtitle: Optional[str]
    text: Optional[List[str]]
    formula_images: Optional[List[Formula]]
    picture: List[str] = []

class Ppt(BaseModel):
    content: List[Slides]

# Cache for analysis results
analysis_cache = {}


# Helper functions
def generate_cache_key(url, student_level):
    return f"{url}_{student_level}"


def clean_data_folder():
    """Ensure data folder contains only paper.pdf"""
    for filename in os.listdir(PDF_DIR):
        if filename != "paper.pdf":
            file_path = os.path.join(PDF_DIR, filename)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception as e:
                print(f"Error deleting {file_path}: {e}")


# Prerequisite parsing function
def parse_prerequisites(text: str) -> dict:
    prerequisites = {}
    sections = re.split(r'### \d+\.\s*', text)
    
    for section in sections[1:]:
        title_match = re.match(r'\*\*(.*?)\*\*', section)
        if title_match:
            title = title_match.group(1).strip()
            sub_content = re.findall(
                r'-\s*\*\*(.*?)\*\*:\s*(.*?)(?=\n\s*-|\Z)', 
                section, 
                re.DOTALL
            )
            prerequisites[title] = [
                f"{item[0].strip()}: {item[1].strip()}" 
                for item in sub_content
            ]
    return prerequisites

@app.get("/student-levels")
async def get_student_levels():
    return {
        "levels": [
            {"id": "1", "name": "PhD Researcher"},
            {"id": "2", "name": "Masters Student"},
            {"id": "3", "name": "Undergraduate Student"}
        ]
    }

def isValidArxivUrl(url: str) -> bool:
    pattern = r'^https?://arxiv\.org/(?:abs|pdf)/\d{4}\.\d+(?:v\d+)?(?:\.pdf)?$'
    return re.match(pattern, url) is not None




@app.post("/analyze/url", response_model=AnalysisResponse)
async def analyze_url(
    url: str = Form(...),
    student_level: str = Form(...),
    settings: Settings = Depends(get_settings)
):
    if not isValidArxivUrl(url):
        raise HTTPException(400, "Invalid arXiv URL format. Expected format: https://arxiv.org/abs/2406.15758 or https://arxiv.org/pdf/2406.15758")

    
    level_map = {
        "1": "phd researcher",
        "2": "masters student",
        "3": "undergraduate student"
    }
    
    if student_level not in level_map:
        raise HTTPException(400, "Invalid student level")
    
    client = Mistral(api_key=settings.MISTRAL_API_KEY)
    # prompt = f"""Analyze this paper and list prerequisites for a {level_map[student_level]}"""
    prompt=f"""Analyze this research paper and provide a comprehensive list of prerequisite topics that a {level_map[student_level]} should be familiar with to fully understand the concepts presented"""
    response = client.chat.complete(
        model=settings.MODEL_NAME,
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": prompt},
                {"type": "document_url", "document_url": url}
            ]
        }]
    )
    json_object = json.dumps(parse_prerequisites(response.choices[0].message.content), indent=4)
    raw_path = Path(config.TEMPLATE_METADATA_DIR) / "prerequisites_dict.json"

    # Ensure the directory exists
    raw_path.parent.mkdir(parents=True, exist_ok=True)

    with open(raw_path, "w") as outfile:
        outfile.write(json_object)
    return {"prerequisites": parse_prerequisites(response.choices[0].message.content)}
    

import shutil

# Add these at the top
PDF_DIR = "data/upload"
os.makedirs(PDF_DIR, exist_ok=True)

# Modified analyze_pdf endpoint
@app.post("/analyze/pdf", response_model=AnalysisResponse)
async def analyze_pdf(
    file: UploadFile = File(...),
    student_level: str = Form(...),
    settings: Settings = Depends(get_settings)
):
    # Clear data folder
    for filename in os.listdir(PDF_DIR):
        file_path = os.path.join(PDF_DIR, filename)
        try:
            if os.path.isfile(file_path):
                os.unlink(file_path)
        except Exception as e:
            print(f"Error deleting {file_path}: {e}")

    # Save new PDF
    clean_data_folder()
    pdf_path = os.path.join(PDF_DIR, "paper.pdf")
    try:
        with open(pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(500, f"Failed to save PDF: {str(e)}")
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are supported")
    
    level_map = {
        "1": "phd researcher",
        "2": "masters student",
        "3": "undergraduate student"
    }
    
    if student_level not in level_map:
        raise HTTPException(status_code=400, detail="Invalid student level. Choose 1, 2, or 3.")
    
    student_level_description = level_map.get(student_level)
    
    try:
        # Set up the Mistral client
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        
        # Read the file content
        file_content = await file.read()
        
        # Upload the PDF to Mistral
        uploaded_pdf = client.files.upload(
            file={
                "file_name": file.filename,
                "content": file_content,
            },
            purpose="ocr"
        )
        
        # Get a signed URL for the uploaded file
        signed_url = client.files.get_signed_url(file_id=uploaded_pdf.id)
        
        # Format the prompt based on student level
        title_prompt = f"""Analyze this research paper and provide a comprehensive list of prerequisite topics that a {student_level_description} should be familiar with to fully understand the concepts presented"""
        
        # Call Mistral API with the signed URL
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": title_prompt
                    },
                    {
                        "type": "document_url",
                        "document_url": signed_url.url
                    }
                ]
            }
        ]
        
        title_chat_response = client.chat.complete(
            model=settings.MODEL_NAME,
            messages=messages
        )
        
        prerequisites = title_chat_response.choices[0].message.content
        parsed_prerequisites = parse_prerequisites(prerequisites)
        json_object = json.dumps(parsed_prerequisites, indent=4)

        raw_path = Path(config.TEMPLATE_METADATA_DIR) / "prerequisites_dict.json"

        # Ensure the directory exists
        raw_path.parent.mkdir(parents=True, exist_ok=True)

        with open(raw_path, "w") as outfile:
            outfile.write(json_object)

        return {"prerequisites": parse_prerequisites(title_chat_response.choices[0].message.content)}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")



@app.post("/extract-template-layout")
async def extract_template_layout(template_name: str = Form("template.pptx")):
    # Use consistent paths and names like the original
    try:
        presentation_path = f"data/upload/{template_name}"
        presentation = Presentation(presentation_path)

        master_slide = presentation.slide_master
        layout_details = []

        for layout in master_slide.slide_layouts:
            layout_info = {"name": layout.name, "placeholders": []}

            for placeholder in layout.placeholders:
                placeholder_info = {
                    "name": placeholder.name,
                    "type": str(placeholder.placeholder_format.type),
                    "index": placeholder.placeholder_format.idx
                }
                layout_info["placeholders"].append(placeholder_info)

            layout_details.append(layout_info)

        # Match original: convert to JSON string and write to file
        layout_details_json = json.dumps(layout_details, indent=4)

        output_dir = Path("data/metadata")
        os.makedirs(output_dir, exist_ok=True)
        output_file = output_dir / "layout_details.json"

        with open(output_file, "w") as file:
            file.write(layout_details_json)

        # Match original return: include both print info and path
        return {
            "message": f"Layout details saved to {output_file}",
            "layout_details": layout_details  # You can also include the raw list if needed
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting layout: {str(e)}")


@app.post("/convert-placeholders")
async def convert_placeholders(layout_extracted_path: str = Form("data/metadata/layout_details.json")):
    try:
        # Load layout details from JSON
        layout_details=load_json(layout_extracted_path)
        # Process each layout
        layout_list = []
        for single_layout_details in layout_details:
            layout_name = single_layout_details["name"]
            layout_dict = {}
            layout_dict["slide_name"] = layout_name
            
            placeholderlist = []
            for placeholder in single_layout_details["placeholders"]:
                placeholderdict = {}
                placeholderdict["name"] = placeholder["name"] + "_" + str(placeholder["index"])
                placeholderlist.append(placeholderdict)

            layout_dict["placeholders"] = placeholderlist
            layout_list.append(layout_dict)

        # Save processed layout to JSON
        raw_path = Path("data/metadata")
        os.makedirs(raw_path, exist_ok=True)
        output_path = raw_path / "processed_layout.json"
        with open(output_path, "w") as f:
            json.dump(layout_details, f, indent=2)
        
        return {"message": "Placeholders converted successfully", "path": str(output_path)}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error converting placeholders: {str(e)}")

@app.get("/use-generated-prerequisite")
async def use_generated_prerequisite(prequisite_json_filename: str = "prerequisites_dict.json"):
    try:
        prequisite_json_path = Path("data/metadata") / prequisite_json_filename
        with open(prequisite_json_path, "r", encoding="utf-8") as f:
            prerequisites_data = json.load(f)

        return {"message": "Prerequisites loaded successfully", "prerequisites": prerequisites_data}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error loading prerequisites: {str(e)}")

@app.post("/ocr-figure-url")
async def ocr_figure_onURL(document_url: str = Form(...), settings: Settings = Depends(get_settings)):
    try:
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        ocr_response = client.ocr.process(
            model=settings.MODEL_NAME_OCR,
            document={
                "type": "document_url",
                "document_url": document_url
            },
            include_image_base64=True
        )
        
        return {"message": "OCR completed successfully", "ocr_response": ocr_response}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error performing OCR: {str(e)}")

@app.post("/ocr-figure-pdf")
async def ocr_figure_pdf(file: UploadFile = File(...), settings: Settings = Depends(get_settings)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "PDF files only")
    
    temp_path = f"data/upload/temp_{file.filename}"
    try:
        with open(temp_path, "wb") as f:
            f.write(await file.read())
        
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        uploaded_pdf = client.files.upload(
            file={
                "file_name": os.path.basename(temp_path),
                "content": open(temp_path, "rb"),
            },
            purpose="ocr"
        )

        signed_url = client.files.get_signed_url(file_id=uploaded_pdf.id)

        ocr_response = client.ocr.process(
            model=settings.MODEL_NAME_OCR,
            document={
                "type": "document_url",
                "document_url": signed_url.url,
            }
        )
        
        return {"message": "OCR completed successfully", "ocr_response": ocr_response}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error performing OCR: {str(e)}")
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

def extract_figure_tag(markdown, image_id):
    """Extracts just 'Figure X' (not full caption) following the image reference."""
    pattern = rf"!\[.*?\]\({re.escape(image_id)}\)\s*[\r\n]+(Figure\s+\d+)"
    match = re.search(pattern, markdown, re.IGNORECASE)
    return match.group(1).strip() if match else None


@app.post("/save-figures")
async def save_figures_from_ocr(ocr_response: dict):
    try:
        output_folder = Path("data/figures")
        os.makedirs(output_folder, exist_ok=True)
        figures_metadata = {}

        for page in ocr_response.get("pages", []):
            page_index = page.get("index")
            markdown = page.get("markdown", "")
            images = page.get("images", [])

            for i, image in enumerate(images):
                image_id = image.get("id", f"figure_{page_index}_{i}.jpeg")
                img_data = image.get("image_base64")

                if img_data:
                    if ',' in img_data:
                        img_data = img_data.split(',', 1)[1]

                    img_bytes = base64.b64decode(img_data)
                    img = Image.open(io.BytesIO(img_bytes))

                    image_filename = os.path.basename(image_id)
                    save_path = output_folder / image_filename

                    img.save(save_path)

                    figure_key = extract_figure_tag(markdown, image_id) or f"Figure_{page_index}_{i}"
                    figures_metadata[figure_key] = str(save_path)

        metadata_path = output_folder / "figures_metadata.json"
        with open(metadata_path, "w") as meta_file:
            json.dump(figures_metadata, meta_file, indent=4)

        return {
            "message": f"Saved {len(figures_metadata)} figures",
            "metadata_path": str(metadata_path),
            "figures_metadata": figures_metadata
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving figures: {str(e)}")

@app.post("/slide-data-gen")
async def slide_data_gen(
    student_level: str = Form(...),
    document_url: str = Form(...),
    num_slides: int = Form(10),
    settings: Settings = Depends(get_settings)
):
    try:
        # Load prerequisites from file
        prerequisites_path = Path("data/metadata/prerequisites_dict.json")
        with open(prerequisites_path, "r") as f:
            prerequisites_dict = json.load(f)
        
        # Initialize Mistral client
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        
        # Create slide prompt
        slide_prompt = f"""Analyze this research paper and create a structured outline for a PowerPoint presentation tailored for a {student_level} audience, where student_level is either 'PhD researcher', 'Master's student', or 'Undergrad student'. Follow these guidelines:

Pictures should only contain figure number, example (figure 2)

ONLY INCLUDE THE TOPICS IN {prerequisites_dict} FOR THE SLIDES

Important: Once a formula is shown do not use it again

CONSIDER Pictures AND FORMULA AS IMAGES

DO NOT KEEP IMAGE/Pictures AND FORMULA IN SAME SLIDE

Total slides should be : {num_slides} **Always follow this, must not be more than this**

Title Slide: Create a concise, engaging title that captures the paper's essence.

This is the first slide and a catchy subtitle. Bullets must be empty in this

Agenda slide: The second slide which contains all subtitles of other slides as its bullet points

Then State the paper's main research question and significance

This acts like the title

Highlight key background information relevant to the {student_level} audience

For each section of the paper:

Create 'Section Overview' slide with:

Section title: bullet points summarizing key concepts in content

Any critical formulas, using LaTeX notation

Followed with in depth explanation slides that only if necessary and also break each point in section overview into one slide explaining in depth:

Explain complex ideas using analogies or visualizations

Break down important formulas step-by-step

Highlight connections to prerequisite knowledge for the {student_level}

Do this for all the chosen topics

Make sure to

Present key findings with supporting data or graphs

Interpret results at a level appropriate for the {student_level}

With great amount of details not missing any key information

Discussion & Implications (1-2 slides):

Outline the paper's main conclusions

Discuss potential applications or future research directions

Relate findings to broader field context for the {student_level}

Key Takeaways slide:

List 3-5 main points to remember, tailored to the {student_level}'s level

Further Reading slide:

Suggest 3-4 related papers or resources appropriate for the {student_level}

For each slide, provide:

A clear, concise headline as the subtitle

Bullet points for main content (3-5 per slide)

Notes on any visuals, charts, or diagrams to include from the paper mention the figure number

Speaker notes with additional context or explanations

Adjust the depth and complexity of the content based on the {student_level}, ensuring the presentation is informative and engaging for the specified audience level.

Make sure to include the mathematical formulas where ever necessary

Give the output in a json format and a dictionary tagging formuala and its name in json
"""        
        # Define the messages for the chat
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": slide_prompt
                    },
                    {
                        "type": "document_url",
                        "document_url": document_url
                    }
                ]
            }
        ]

        chat_response = client.chat.parse(
            model=settings.MODEL_NAME,
            messages=messages,
            response_format=Ppt
        )

        slides = chat_response.choices[0].message.content
        raw_path = Path("data/metadata") / "slides_data.json"
        with open(raw_path, "w") as f:
            json.dump(slides, f, indent=2)
        
        return {"message": "Slide data generated successfully", "path": str(raw_path)}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating slide data: {str(e)}")

def render_latex_to_image(formula, name="latex"):
    try:
        fig, ax = plt.subplots()
        fig.patch.set_visible(False)
        image_path = Path("data/formulas") / f"{name.replace(' ', '_')}.png"
        ax.axis('off')
        ax.text(0.5, 0.5, f"{formula}", fontsize=30, ha='center', va='center')
        image_path.parent.mkdir(parents=True, exist_ok=True)
        plt.savefig(image_path, bbox_inches='tight', transparent=True)
        plt.close()
        return str(image_path)
    except:
        try:
            fig, ax = plt.subplots()
            fig.patch.set_visible(False)
            image_path = Path("data/formulas") / f"{name.replace(' ', '_')}.png"
            ax.axis('off')
            ax.text(0.5, 0.5, f"${formula}$", fontsize=30, ha='center', va='center')
            image_path.parent.mkdir(parents=True, exist_ok=True)
            plt.savefig(image_path, bbox_inches='tight', transparent=True)
            plt.close()
            return str(image_path)
        except Exception as e:
            print(f"❌ Error rendering formula '{name}': {str(e)}")
            return ""


@app.post("/process-slides-data")
async def process_slides_data():
    try:
        print("\n🔥 Starting /process-slides-data")

        input_slides_path = Path("data/metadata/slides_data.json")
        output_slides_path = Path("data/metadata/updated_slides_data.json")
        figures_metadata_path=Path(config.FIGURES_DIR) / "figures_metadata.json"
        # figures_metadata_path = Path("data/figures/figures_metadata.json")

        if not input_slides_path.exists():
            raise HTTPException(status_code=400, detail="Missing slides_data.json")
        # if not figures_metadata_path.exists():
        #     raise HTTPException(status_code=400, detail="Missing figures_metadata.json")

        # with open(input_slides_path, "r") as file:
        #     slides_data = json.load(file)
        slides_data=load_json(input_slides_path)

        # with open(figures_metadata_path, "r") as file:
        #     image_data = json.load(file)
        image_data=load_json(figures_metadata_path)

        # Process formulas
        for slide in slides_data.get("content", []):
            if "formula_images" in slide and slide["formula_images"]:
                for item in slide["formula_images"]:
                    if isinstance(item, dict) and "formula" in item and "name" in item:
                        formula = item["formula"].strip('$')
                        name = item["name"]
                        image_filename = render_latex_to_image(formula, name)
                        item["formula"] = image_filename

        # Process pictures
        for slide in slides_data.get("content", []):
            if "picture" in slide and slide["picture"]:
                for i in range(len(slide["picture"])):
                    key = slide["picture"][i].capitalize()
                    if key in image_data:
                        slide["picture"][i] = image_data[key]
                    else:
                        print(f"⚠️ Warning: Missing image data for key '{key}'")
                        slide["picture"][i] = ""

        # Save the updated JSON
        with open(output_slides_path, "w") as file:
            json.dump(slides_data, file, indent=4)

        print("✅ Formulas and pictures processed successfully.")
        return {"message": "Slides processed and updated JSON saved.", "path": str(output_slides_path)}

    except HTTPException as he:
        print(f"🚨 HTTP Error {he.status_code}: {he.detail}")
        raise he
    except Exception as e:
        print(f"🚨 Unexpected Error: {type(e).__name__}: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")
    
@app.post("/execution-agent-parsing")
async def execution_agent_parsing(
    template_name: str = Form(...),
    settings: Settings = Depends(get_settings)
):
    try:
        print(f"📌 Execution Agent using template: {template_name}")
        # Check for template with case-insensitive matching
        upload_dir = Path("data/upload")
        template_path = None
        
        # Try exact match first
        exact_match = upload_dir / template_name
        if exact_match.exists():
            template_path = exact_match
        else:
            # Try case-insensitive matching
            for file in upload_dir.glob("*.pptx"):
                if file.name.lower() == template_name.lower():
                    template_path = file
                    print(f"📌 Found template with case-insensitive match: {file.name}")
                    break
        
        if not template_path:
            raise HTTPException(
                status_code=400, 
                detail=f"Template file {template_name} not found. Available templates: {', '.join([f.name for f in upload_dir.glob('*.pptx')])}"
            )
            
        template_dir = template_path
        json_dir = Path("data/metadata")
        slides_data_path = json_dir / "updated_slides_data.json"
        layout_data_path = json_dir / "processed_layout.json"
        output_path = json_dir / "execution_agent.json"

        presentation = Presentation(template_dir)
        distinct_layout = [layout.name for layout in presentation.slide_master.slide_layouts]

        slides_data = load_json(slides_data_path)
        slides_layout = load_json(layout_data_path)

        query=f"""
        Understand the current slides data as provided:
        {slides_data}

        Now create a new json with layouts chosen from {distinct_layout} for each of the slide. The json contains the placeholders.name_placeholders.index from {slides_layout} as the key and the content as the values.

        Using the fields in the slides data namely title, subtitle, text, formula_images and picture to the placeholder of the layout selected for the slide.

        Give new titles for each slide and remove the subtitle

        Treat both formulas and pictures as images when assigning the layout

        If the content has bullet points, then be creative in choosing layouts which contain texts and title type

        ## And do not loose the content ##

        ##final output must in JSON only
        """
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        response = client.agents.complete(
            agent_id=settings.EXECUTION_AGENT_ID,
            messages=[{"role": "user", "content": query}]
        )
        agent_result = response.choices[0].message.content
        print("Got the agent Response: ",agent_result)
        def extract_json(text):
            match = re.search(r'```json(.*?)```', text, re.DOTALL)
            if match:
                json_str = match.group(1).strip()
                try:
                    return json.loads(json_str)
                except json.JSONDecodeError as e:
                    print(f"JSON Decoding Error: {e}")
                    return None
            else:
                print("No JSON block found in the agent output.")
                return None


        execution_agent_json = extract_json(agent_result)

        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(execution_agent_json, f, indent=2)

        return {"message": "Execution agent parsing completed", "path": str(output_path)}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error in execution agent parsing: {str(e)}")

import traceback
    
@app.post("/generate-presentation")
async def generate_presentation_from_execution_json(
    template_name: str = Form(...),
    execution_json_filename: str = Form("execution_agent.json"),
    output_ppt_filename: str = Form("modified_presentation.pptx"),
    processed_layout_filename: str = Form("processed_layout.json")
):
    print(f"📌 Generating presentation with template: {template_name}")
    try:
        # Prepare directories
        output_dir = Path("data/output")
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Find the template file with case-insensitive matching
        upload_dir = Path("data/upload")
        template_path = None
        
        # Try exact match first
        exact_match = upload_dir / template_name
        if exact_match.exists():
            template_path = exact_match
            print(f"📌 Using exact template match: {template_name}")
        else:
            # Try case-insensitive matching
            for file in upload_dir.glob("*.pptx"):
                if file.name.lower() == template_name.lower():
                    template_path = file
                    print(f"📌 Found template with case-insensitive match: {file.name}")
                    break
        
        if not template_path:
            raise HTTPException(
                status_code=400, 
                detail=f"Template file {template_name} not found. Available templates: {', '.join([f.name for f in upload_dir.glob('*.pptx')])}"
            )
            
        # Set up paths
        template_dir = template_path
        json_dir = Path("data/metadata")
        execution_agent_json = json_dir / execution_json_filename
        output_ppt_path = Path("data/output") / output_ppt_filename

        # Load the presentation and JSON data
        prs = Presentation(template_dir)
        print(f"📄 Presentation loaded from {template_dir}")
        
        # Get the initial count of slides for later deletion
        initial_slide_count = len(prs.slides)
        print(f"📊 Template has {initial_slide_count} initial slides")
        
        # Load the JSON data
        json_data = load_json(execution_agent_json)
        print(f"📊 Loaded execution agent JSON with {len(json_data.get('slides', []))} slides")

        # Define layout finder function
        def get_layout_by_name(prs_obj, layout_name):
            for layout in prs_obj.slide_layouts:
                if layout.name == layout_name:
                    return layout
            return None

        # For debugging
        layout_names = [layout.name for layout in prs.slide_layouts]
        print(f"📋 Available layouts: {layout_names}")

        # Create slides from JSON data
        print("🔄 Creating slides from JSON data...")
        for slide_idx, slide_data in enumerate(json_data.get("slides", []), 1):
            layout_name = slide_data.get("slide_name")
            layout = get_layout_by_name(prs, layout_name)
            
            if not layout:
                print(f"⚠️ Layout '{layout_name}' not found in the template. Skipping slide {slide_idx}.")
                continue
                
            # Add the slide
            slide = prs.slides.add_slide(layout)
            print(f"✅ Added slide {slide_idx} with layout '{layout_name}'")
            
            # Fill placeholders
            for placeholder_name, content in slide_data.get("placeholders", {}).items():
                # Parse the placeholder name to get components
                name_parts = placeholder_name.split("_")
                name = name_parts[0] if name_parts else placeholder_name
                index = name_parts[-1] if len(name_parts) > 1 else "0"
                
                try:
                    # Try to find the placeholder by index
                    placeholder_found = False
                    idx = int(index)
                    for shape in slide.shapes:
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == idx:
                            placeholder_found = True
                            
                            # Handle different content types
                            if shape.has_text_frame:
                                # Convert content to string if it's a list
                                if isinstance(content, str):
                                    # Check if this is a formula path
                                    if ("formulas" in content or "Formula" in content) and content.endswith(".png"):
                                        try:
                                            # This is a formula image, insert it
                                            formula_path = content if content.startswith("data/") else f"data/formulas/{content}"
                                            print(f"🔤 Inserting formula image at {formula_path}")
                                            shape.insert_picture(formula_path)
                                        except Exception as e:
                                            print(f"⚠️ Error inserting formula image: {e}")
                                            shape.text_frame.text = ""  # Don't show the path
                                    else:
                                        shape.text_frame.text = content
                                elif isinstance(content, list):
                                    shape.text_frame.text = "\n".join([str(item) for item in content])
                                else:
                                    shape.text_frame.text = str(content)
                                    
                            # Handle picture placeholders
                            if name.startswith("Picture") or "Picture" in shape.name:
                                try:
                                    # Determine the path
                                    pic_path = content
                                    if isinstance(content, list) and content:
                                        pic_path = content[0]
                                        
                                    # Make sure it has the full path
                                    if isinstance(pic_path, str):
                                        if not pic_path.startswith("data/"):
                                            if "Formula" in pic_path:
                                                pic_path = f"data/formulas/{pic_path}"
                                            elif "figure" in pic_path.lower() or "Figure" in pic_path:
                                                pic_path = f"data/figures/{pic_path}"
                                                
                                        print(f"🖼️ Inserting picture: {pic_path}")
                                        shape.insert_picture(pic_path)
                                except Exception as e:
                                    print(f"⚠️ Error inserting picture: {e}")
                            
                    if not placeholder_found:
                        print(f"⚠️ Placeholder with index {idx} not found in slide {slide_idx}")
                        
                except ValueError:
                    # If the index isn't a number, try to find by name
                    placeholder_found = False
                    for shape in slide.shapes:
                        if hasattr(shape, 'name') and name in shape.name:
                            placeholder_found = True
                            if shape.has_text_frame:
                                if isinstance(content, str):
                                    # Check if this is a formula path
                                    if ("formulas" in content or "Formula" in content) and content.endswith(".png"):
                                        try:
                                            # This is a formula image, insert it
                                            formula_path = content if content.startswith("data/") else f"data/formulas/{content}"
                                            print(f"🔤 Inserting formula image at {formula_path}")
                                            shape.insert_picture(formula_path)
                                        except Exception as e:
                                            print(f"⚠️ Error inserting formula image: {e}")
                                            shape.text_frame.text = ""  # Don't show the path
                                    else:
                                        shape.text_frame.text = content
                                elif isinstance(content, list):
                                    shape.text_frame.text = "\n".join([str(item) for item in content])
                                else:
                                    shape.text_frame.text = str(content)
                            
                            # Handle picture placeholders by name
                            if name.startswith("Picture") or "Picture" in shape.name:
                                try:
                                    # Determine the path
                                    pic_path = content
                                    if isinstance(content, list) and content:
                                        pic_path = content[0]
                                        
                                    # Make sure it has the full path
                                    if isinstance(pic_path, str):
                                        if not pic_path.startswith("data/"):
                                            if "Formula" in pic_path:
                                                pic_path = f"data/formulas/{pic_path}"
                                            elif "figure" in pic_path.lower() or "Figure" in pic_path:
                                                pic_path = f"data/figures/{pic_path}"
                                                
                                        print(f"🖼️ Inserting picture by name: {pic_path}")
                                        shape.insert_picture(pic_path)
                                except Exception as e:
                                    print(f"⚠️ Error inserting picture by name: {e}")
                    
                    if not placeholder_found:
                        print(f"⚠️ Placeholder with name {name} not found in slide {slide_idx}")

        # Save the presentation with all slides included
        temp_ppt_path = "temp.pptx"
        prs.save(temp_ppt_path)
        print(f"💾 Saved temporary presentation to {temp_ppt_path}")

        print(f"🛠️ Now handling slide deletion following the original code's approach")
        # Follow the exact approach from the notebook code:
        # Load the presentation
        prs = Presentation(temp_ppt_path)
        
        # Get the total number of slides
        total_slides = len(prs.slides)
        print(f"📊 Total slides in presentation: {total_slides}")
        
        # Use exactly 13 as in the original code, or fewer if there are fewer slides
        slides_to_remove = min(13, total_slides)
        print(f"🧹 Removing {slides_to_remove} slides from the beginning")
        
        # Remove the slides from the beginning
        for _ in range(slides_to_remove):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        # Save the final presentation
        prs.save(str(output_ppt_path).replace("\\","/"))
        print(f"✅ Final presentation saved to {output_ppt_path}")
        
        # Clean up temp file
        os.remove(temp_ppt_path)

        return {"message": "Presentation generated successfully", "path": str(output_ppt_path)}
    
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error generating presentation: {str(e)}")


@app.get("/download-presentation")
async def download_presentation(filename: str = "modified_presentation.pptx"):
    file_path = Path("data/output") / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Presentation file not found")
    return FileResponse(file_path, filename=filename)

# Add this new endpoint for cleaning up data
@app.post("/cleanup-data")
async def cleanup_data():
    """Clean up figures and formulas directories after successful generation"""
    try:
        # Clear the figures directory
        figures_dir = Path("data/figures")
        if figures_dir.exists():
            for item in figures_dir.glob("*"):
                if item.is_file():
                    item.unlink()
                elif item.is_dir():
                    shutil.rmtree(item)
        
        # Clear the formulas directory
        formulas_dir = Path("data/formulas")
        if formulas_dir.exists():
            for item in formulas_dir.glob("*"):
                if item.is_file():
                    item.unlink()
                elif item.is_dir():
                    shutil.rmtree(item)
                    
        return {"message": "Successfully cleaned up figures and formulas data"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error cleaning up data: {str(e)}")


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
