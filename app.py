from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from mistralai import Mistral
import json
import re
import os
import hashlib
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
from pydantic_settings import BaseSettings
from functools import lru_cache
import uvicorn
from pathlib import Path
from pptx import Presentation
from PIL import Image
import base64
import io
import concurrent.futures
import matplotlib
matplotlib.use('Agg')  # Use the 'Agg' backend which doesn't require a GUI
import matplotlib.pyplot as plt
from sympy import preview
import shutil
import traceback
import time
import asyncio
from config import config
from utils import load_json, save_json, extract_json
from mistralai.client import MistralClient
# from mistralai.models.chat_completion import ChatMessage
import random
from pydub import AudioSegment
from melo.api import TTS
import json
from podcast_api import generate_podcast

# Configuration settings
class Settings(BaseSettings):
    MISTRAL_API_KEY: str
    MODEL_NAME: str = "mistral-small-latest"
    EXECUTION_AGENT_ID: Optional[str] = None
    MODEL_NAME_OCR: str = "mistral-ocr-latest"
    ENHANCE_AGENT_ID: Optional[str] = None
    class Config:
        env_file = ".env"
        extra = 'ignore'

@lru_cache
def get_settings():
    return Settings()

# Create FastAPI app
app = FastAPI(title="Research Paper Prerequisites API")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create necessary directories
os.makedirs("data/upload", exist_ok=True)
os.makedirs("data/figures", exist_ok=True)
os.makedirs("data/formulas", exist_ok=True)
os.makedirs("data/metadata", exist_ok=True)
os.makedirs("data/output", exist_ok=True)
os.makedirs("podcast", exist_ok=True)
os.makedirs("images", exist_ok=True)

# Mount the images directory
app.mount("/images", StaticFiles(directory="images"), name="images")

# Define models
class AnalysisResponse(BaseModel):
    prerequisites: dict

class Formula(BaseModel):
    formula: str
    name: str

class Slides(BaseModel):
    title: str
    subtitle: Optional[str]
    text: Optional[List[str]]
    formula_images: Optional[List[Formula]]
    picture: List[str] = []

class Ppt(BaseModel):
    content: List[Slides]


class Slide(BaseModel):
    slide_name: str
    placeholders: Dict[str, str]

class PresentationData(BaseModel):
    slides: List[Slide]


# Cache for analysis results
analysis_cache = {}
formula_cache = {}

# Directory for PDFs
PDF_DIR = "data/upload"
os.makedirs(PDF_DIR, exist_ok=True)

# Helper functions
def isValidArxivUrl(url: str) -> bool:
    """Check if URL is a valid arXiv URL"""
    pattern = r'^https?://arxiv\.org/(?:abs|pdf)/\d{4}\.\d+(?:v\d+)?(?:\.pdf)?$'
    return re.match(pattern, url) is not None

def generate_cache_key(url, student_level):
    """Generate cache key for analysis results"""
    return f"{url}_{student_level}"

def clean_data_folder():
    """Ensure data folder contains only paper.pdf"""
    for filename in os.listdir(PDF_DIR):
        if filename != "paper.pdf":
            file_path = os.path.join(PDF_DIR, filename)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception as e:
                print(f"Error deleting {file_path}: {e}")

def extract_figure_tag(markdown, image_id):
    """Extracts just 'Figure X' (not full caption) following the image reference."""
    pattern = rf"!\[.*?\]\({re.escape(image_id)}\)\s*[\r\n]+(Figure\s+\d+)"
    match = re.search(pattern, markdown, re.IGNORECASE)
    return match.group(1).strip() if match else None

# Prerequisite parsing function
def parse_prerequisites(text: str) -> dict:
    """Parse prerequisites from Mistral response text"""
    prerequisites = {}
    sections = re.split(r'### \d+\.\s*', text)
    
    for section in sections[1:]:
        title_match = re.match(r'\*\*(.*?)\*\*', section)
        if title_match:
            title = title_match.group(1).strip()
            sub_content = re.findall(
                r'-\s*\*\*(.*?)\*\*:\s*(.*?)(?=\n\s*-|\Z)', 
                section, 
                re.DOTALL
            )
            prerequisites[title] = [
                f"{item[0].strip()}: {item[1].strip()}" 
                for item in sub_content
            ]
    return prerequisites

def render_latex_to_image(formula: str, name: str="latex", dpi: int=200) -> str:
    """
    Generate high-quality formula PNG using LaTeX + dvipng.

    Args:
    - formula: Raw LaTeX math content, without enclosing $...$.
    - name: Identifier for the file name; the output will be saved as data/formulas/{name}.png
    - dpi: Resolution (dots per inch) for the output PNG.

    Returns:
    - Path to the generated PNG file.
    """
    out_dir = Path("data/formulas")
    out_dir.mkdir(parents=True, exist_ok=True)
    output_path = out_dir / f"{name.replace(' ', '_')}.png"

    # Sympy.preview will:
    #   1) Write a minimal standalone LaTeX document in a temporary directory
    #   2) Compile it to DVI/PDF using pdflatex
    #   3) Use dvipng to convert the output into a PNG image

    try:
        preview(
            f"${formula}$",
            output='png',
            viewer='file',
            filename=str(output_path),
            dvioptions=[f"-D{dpi}", "-Ttight"]
        )
        return str(output_path)
    except Exception as e:
        print(f"❌ render_latex_to_image failed，back to Matplotlib: {e}")

    # Matplotlib + usetex
    try:
        plt.rcParams.update({"text.usetex": True})
        fig = plt.figure(figsize=(0.01, 0.01))
        fig.patch.set_visible(False)
        fig.text(0.5, 0.5, f"${formula}$", ha='center', va='center', fontsize=20)
        plt.savefig(output_path, dpi=dpi, bbox_inches='tight', pad_inches=0.1, transparent=True)
        plt.close(fig)
        return str(output_path)
    except Exception as e2:
        print(f"⚠️ Matplotlib also failed：{e2}")
        return ""

# # Function to render LaTeX formulas as images
# def render_latex_to_image(formula, name="latex"):
#     """Render LaTeX formula to image using matplotlib"""
#     # Check cache first
#     cache_key = f"{formula}_{name}"
#     if cache_key in formula_cache:
#         return formula_cache[cache_key]
    
#     try:
#         start_time = time.time()
#         fig, ax = plt.subplots()
#         fig.patch.set_visible(False)
#         image_path = Path("data/formulas") / f"{name.replace(' ', '_')}.png"
#         ax.axis('off')
#         ax.text(0.5, 0.5, f"${formula}$", fontsize=30, ha='center', va='center')
#         image_path.parent.mkdir(parents=True, exist_ok=True)
#         plt.savefig(image_path, bbox_inches='tight', transparent=True, dpi=150)
#         plt.close()
        
#         render_time = time.time() - start_time
#         print(f"⚡ Formula '{name}' rendered in {render_time:.2f}s with $ escape")
        
#         # Cache the result
#         formula_cache[cache_key] = str(image_path)
#         return str(image_path)
        
#     except:
#         try:
#             start_time = time.time()
#             fig, ax = plt.subplots()
#             fig.patch.set_visible(False)
#             image_path = Path("data/formulas") / f"{name.replace(' ', '_')}.png"
#             ax.axis('off')
#             ax.text(0.5, 0.5, f"{formula}", fontsize=30, ha='center', va='center')
#             image_path.parent.mkdir(parents=True, exist_ok=True)
#             plt.savefig(image_path, bbox_inches='tight', transparent=True, dpi=150)
#             plt.close()
            
#             render_time = time.time() - start_time
#             print(f"⚡ Formula '{name}' rendered in {render_time:.2f}s")
            
#             # Cache the result
#             formula_cache[cache_key] = str(image_path)
#             return str(image_path)
#         except Exception as e:
#             print(f"❌ Error rendering formula '{name}': {str(e)}")
#             return ""

# Directory cleaning function
def clean_directories(directories=None):
    """
    Clean specified directories by removing all files within them.
    If directories is None, clean all data directories except upload and images.
    """
    if directories is None:
        directories = [
            "data/figures", 
            "data/formulas", 
            "data/metadata",
            "podcast"
        ]
    
    start_time = time.time()
    print(f"🧹 Cleaning directories: {directories}")
    
    # Ensure images directory is not accidentally cleaned
    directories = [dir for dir in directories if dir != "images"]
    
    for directory in directories:
        dir_path = Path(directory)
        if dir_path.exists():
            for filename in os.listdir(dir_path):
                file_path = dir_path / filename
                try:
                    if file_path.is_file():
                        os.unlink(file_path)
                    elif file_path.is_dir():
                        shutil.rmtree(file_path)
                except Exception as e:
                    print(f"❌ Error cleaning {file_path}: {e}")
    
    elapsed = time.time() - start_time
    print(f"✅ Directory cleaning completed in {elapsed:.2f}s")

# Parallel formula processing
def process_formulas_parallel(slides_data):
    """Process all formulas in slides data in parallel for better performance"""
    formula_tasks = []
    
    # Collect all formulas that need rendering
    for slide in slides_data.get("content", []):
        if "formula_images" in slide and slide["formula_images"]:
            for item in slide["formula_images"]:
                if isinstance(item, dict) and "formula" in item and "name" in item:
                    formula = item["formula"].strip('$')
                    name = item["name"]
                    formula_tasks.append((formula, name, item))
    
    # Render formulas in parallel
    if formula_tasks:
        print(f"🧮 Processing {len(formula_tasks)} formulas")
        for formula, name, item in formula_tasks[:3]:  # Print first 3 for debugging
            print(f"  - Formula: '{name}' -> '{formula[:30]}...'")
        start_time = time.time()
        print(f"🧮 Processing {len(formula_tasks)} formulas in parallel")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(8, len(formula_tasks))) as executor:
            futures = {executor.submit(render_latex_to_image, formula, name): (formula, name, item) 
                    for formula, name, item in formula_tasks}
            
            for future in concurrent.futures.as_completed(futures):
                _, _, item = futures[future]
                try:
                    image_filename = future.result()
                    item["formula"] = image_filename
                except Exception as e:
                    print(f"❌ Error processing formula: {e}")
        
        elapsed = time.time() - start_time
        print(f"✅ Formula processing completed in {elapsed:.2f}s")
    
    return slides_data 

# API Endpoints

@app.get("/student-levels")
async def get_student_levels():
    """Get available student levels"""
    return {
        "levels": [
            {"id": "1", "name": "PhD Researcher"},
            {"id": "2", "name": "Masters Student"},
            {"id": "3", "name": "Undergraduate Student"}
        ]
    }

@app.post("/analyze/url", response_model=AnalysisResponse)
async def analyze_url(
    url: str = Form(...),
    student_level: str = Form(...),
    settings: Settings = Depends(get_settings)
):
    """Analyze paper from URL and extract prerequisites"""
    # Clean directories before starting new analysis
    clean_directories()
    
    if not isValidArxivUrl(url):
        raise HTTPException(400, "Invalid arXiv URL format. Expected format: https://arxiv.org/abs/2406.15758 or https://arxiv.org/pdf/2406.15758")
    
    level_map = {
        "1": "phd researcher",
        "2": "masters student",
        "3": "undergraduate student"
    }
    
    if student_level not in level_map:
        raise HTTPException(400, "Invalid student level")
    
    # Start API call timer for performance monitoring
    start_time = time.time()
    
    try:
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        prompt = f"""Analyze this research paper and provide a comprehensive list of prerequisite topics that a {level_map[student_level]} should be familiar with to fully understand the concepts presented"""
        
        # Make the API call
        response = client.chat.complete(
            model=settings.MODEL_NAME,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "document_url", "document_url": url}
                ]
            }]
        )
        
        # Log performance metrics
        api_time = time.time() - start_time
        print(f"📊 API call completed in {api_time:.2f}s")
        
        # Parse the response and save results
        parsed_prerequisites = parse_prerequisites(response.choices[0].message.content)
        json_object = json.dumps(parsed_prerequisites, indent=4)
        raw_path = Path(config.TEMPLATE_METADATA_DIR) / "prerequisites_dict.json"

        # Ensure the directory exists
        raw_path.parent.mkdir(parents=True, exist_ok=True)

        with open(raw_path, "w") as outfile:
            outfile.write(json_object)
            
        total_time = time.time() - start_time
        print(f"✅ Total processing time: {total_time:.2f}s")
        
        return {"prerequisites": parsed_prerequisites}
    
    except Exception as e:
        print(f"❌ Error analyzing URL: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error analyzing paper: {str(e)}")

@app.post("/analyze/pdf", response_model=AnalysisResponse)
async def analyze_pdf(
    file: UploadFile = File(...),
    student_level: str = Form(...),
    settings: Settings = Depends(get_settings)
):
    """Analyze paper from uploaded PDF and extract prerequisites"""
    # Clean directories before starting new analysis
    clean_directories()
    
    # Validate file
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are supported")
    
    # Save new PDF
    pdf_path = os.path.join(PDF_DIR, "paper.pdf")
    try:
        with open(pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(500, f"Failed to save PDF: {str(e)}")
    
    # Validate student level
    level_map = {
        "1": "phd researcher",
        "2": "masters student",
        "3": "undergraduate student"
    }
    
    if student_level not in level_map:
        raise HTTPException(status_code=400, detail="Invalid student level. Choose 1, 2, or 3.")
    
    student_level_description = level_map.get(student_level)
    
    # Start timer for performance monitoring
    start_time = time.time()
    
    try:
        # Set up the Mistral client
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        
        # Read the file content
        file.file.seek(0)  # Reset file pointer
        file_content = await file.read()
        
        # Upload the PDF to Mistral
        upload_start = time.time()
        uploaded_pdf = client.files.upload(
            file={
                "file_name": file.filename,
                "content": file_content,
            },
            purpose="ocr"
        )
        upload_time = time.time() - upload_start
        print(f"📤 PDF upload completed in {upload_time:.2f}s")
        
        # Get a signed URL for the uploaded file
        signed_url = client.files.get_signed_url(file_id=uploaded_pdf.id)
        
        # Format the prompt based on student level
        title_prompt = f"""Analyze this research paper and provide a comprehensive list of prerequisite topics that a {student_level_description} should be familiar with to fully understand the concepts presented"""
        
        # Call Mistral API with the signed URL
        api_start = time.time()
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": title_prompt
                    },
                    {
                        "type": "document_url",
                        "document_url": signed_url.url
                    }
                ]
            }
        ]
        
        title_chat_response = client.chat.complete(
            model=settings.MODEL_NAME,
            messages=messages
        )
        api_time = time.time() - api_start
        print(f"📊 API analysis completed in {api_time:.2f}s")
        
        # Parse and save results
        parse_start = time.time()
        parsed_prerequisites = parse_prerequisites(title_chat_response.choices[0].message.content)
        json_object = json.dumps(parsed_prerequisites, indent=4)

        raw_path = Path(config.TEMPLATE_METADATA_DIR) / "prerequisites_dict.json"
        raw_path.parent.mkdir(parents=True, exist_ok=True)

        with open(raw_path, "w") as outfile:
            outfile.write(json_object)
            
        parse_time = time.time() - parse_start
        print(f"🔍 Parsing and saving completed in {parse_time:.2f}s")

        total_time = time.time() - start_time
        print(f"✅ Total processing time: {total_time:.2f}s")

        return {"prerequisites": parsed_prerequisites}
    
    except Exception as e:
        print(f"❌ Error analyzing PDF: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error analyzing paper: {str(e)}")

@app.get("/use-generated-prerequisite")
async def use_generated_prerequisite(prequisite_json_filename: str = "prerequisites_dict.json"):
    """Load prerequisites from saved JSON file"""
    try:
        prequisite_json_path = Path("data/metadata") / prequisite_json_filename
        prerequisites_data = load_json(prequisite_json_path)
        return {"message": "Prerequisites loaded successfully", "prerequisites": prerequisites_data}
    
    except Exception as e:
        print(f"❌ Error loading prerequisites: {e}")
        raise HTTPException(status_code=500, detail=f"Error loading prerequisites: {str(e)}")

@app.post("/extract-template-layout")
async def extract_template_layout(template_name: str = Form("template.pptx")):
    """Extract layout from PowerPoint template"""
    try:
        # Check for template with case-insensitive matching
        upload_dir = Path("data/upload")
        template_path = None
        
        # Try exact match first
        exact_match = upload_dir / template_name
        if exact_match.exists():
            template_path = exact_match
            print(f"📌 Using exact template match: {template_name}")
        else:
            # Try case-insensitive matching
            for file in upload_dir.glob("*.pptx"):
                if file.name.lower() == template_name.lower():
                    template_path = file
                    print(f"📌 Found template with case-insensitive match: {file.name}")
                    break
        
        if not template_path:
            raise HTTPException(
                status_code=400, 
                detail=f"Template file {template_name} not found. Available templates: {', '.join([f.name for f in upload_dir.glob('*.pptx')])}"
            )
        
        presentation = Presentation(template_path)
        master_slide = presentation.slide_master
        layout_details = []

        for layout in master_slide.slide_layouts:
            layout_info = {"name": layout.name, "placeholders": []}

            for placeholder in layout.placeholders:
                placeholder_info = {
                    "name": placeholder.name,
                    "type": str(placeholder.placeholder_format.type),
                    "index": placeholder.placeholder_format.idx
                }
                layout_info["placeholders"].append(placeholder_info)

            layout_details.append(layout_info)

        # Convert to JSON string and write to file
        layout_details_json = json.dumps(layout_details, indent=4)

        output_dir = Path("data/metadata")
        os.makedirs(output_dir, exist_ok=True)
        output_file = output_dir / "layout_details.json"

        with open(output_file, "w") as file:
            file.write(layout_details_json)

        return {
            "message": f"Layout details saved to {output_file}",
            "layout_details": layout_details
        }

    except Exception as e:
        print(f"❌ Error extracting layout: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error extracting layout: {str(e)}")

@app.post("/convert-placeholders")
async def convert_placeholders(layout_extracted_path: str = Form("data/metadata/layout_details.json")):
    """Convert layout details to placeholder format"""
    try:
        # Load layout details from JSON
        layout_details = load_json(layout_extracted_path)
        
        # Process each layout
        layout_list = []
        for single_layout_details in layout_details:
            layout_name = single_layout_details["name"]
            layout_dict = {}
            layout_dict["slide_name"] = layout_name
            
            placeholderlist = []
            for placeholder in single_layout_details["placeholders"]:
                placeholderdict = {}
                placeholderdict["name"] = placeholder["name"] + "_" + str(placeholder["index"])
                placeholderlist.append(placeholderdict)

            layout_dict["placeholders"] = placeholderlist
            layout_list.append(layout_dict)

        # Save processed layout to JSON
        raw_path = Path(config.TEMPLATE_METADATA_DIR)
        os.makedirs(raw_path, exist_ok=True)
        output_path = raw_path / "processed_layout.json"
        save_json(layout_list, output_path)
        
        return {"message": "Placeholders converted successfully", "path": str(output_path)}
    
    except Exception as e:
        print(f"❌ Error converting placeholders: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error converting placeholders: {str(e)}")

@app.post("/ocr-figure-url")
async def ocr_figure_onURL(document_url: str = Form(...), settings: Settings = Depends(get_settings)):
    """Perform OCR on a document URL to extract figures"""
    try:
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        print(f"🔍 Performing OCR on URL: {document_url}")
        
        ocr_response = client.ocr.process(
            model=settings.MODEL_NAME_OCR,
            document={
                "type": "document_url",
                "document_url": document_url
            },
            include_image_base64=True
        )
        
        # print(f"✅ OCR completed, extracted {len(ocr_response.get('pages', []))} pages")
        return {"message": "OCR completed successfully", "ocr_response": ocr_response}
    
    except Exception as e:
        print(f"❌ Error performing OCR: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error performing OCR: {str(e)}")

@app.post("/ocr-figure-pdf")
async def ocr_figure_pdf(file: UploadFile = File(...), settings: Settings = Depends(get_settings)):
    """Perform OCR on an uploaded PDF to extract figures"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "PDF files only")
    
    temp_path = f"data/upload/temp_{file.filename}"
    try:
        with open(temp_path, "wb") as f:
            f.write(await file.read())
        
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        uploaded_pdf = client.files.upload(
            file={
                "file_name": os.path.basename(temp_path),
                "content": open(temp_path, "rb"),
            },
            purpose="ocr"
        )

        signed_url = client.files.get_signed_url(file_id=uploaded_pdf.id)

        ocr_response = client.ocr.process(
            model=settings.MODEL_NAME_OCR,
            document={
                "type": "document_url",
                "document_url": signed_url.url,
            }
        )
        
        return {"message": "OCR completed successfully", "ocr_response": ocr_response}
    
    except Exception as e:
        print(f"❌ Error performing OCR: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error performing OCR: {str(e)}")
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/save-figures")
async def save_figures_from_ocr(ocr_response: dict):
    """Save figures extracted from OCR response"""
    try:
        output_folder = Path("data/figures")
        os.makedirs(output_folder, exist_ok=True)
        figures_metadata = {}

        for page in ocr_response.get("pages", []):
            page_index = page.get("index")
            markdown = page.get("markdown", "")
            images = page.get("images", [])

            for i, image in enumerate(images):
                image_id = image.get("id", f"figure_{page_index}_{i}.jpeg")
                img_data = image.get("image_base64")

                if img_data:
                    if ',' in img_data:
                        img_data = img_data.split(',', 1)[1]

                    img_bytes = base64.b64decode(img_data)
                    img = Image.open(io.BytesIO(img_bytes))

                    image_filename = os.path.basename(image_id)
                    save_path = output_folder / image_filename

                    img.save(save_path)

                    figure_key = extract_figure_tag(markdown, image_id) or f"Figure_{page_index}_{i}"
                    figures_metadata[figure_key] = str(save_path)

        metadata_path = output_folder / "figures_metadata.json"
        with open(metadata_path, "w") as meta_file:
            json.dump(figures_metadata, meta_file, indent=4)

        print(f"✅ Saved {len(figures_metadata)} figures")
        return {
            "message": f"Saved {len(figures_metadata)} figures",
            "metadata_path": str(metadata_path),
            "figures_metadata": figures_metadata
        }
    
    except Exception as e:
        print(f"❌ Error saving figures: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error saving figures: {str(e)}")

@app.post("/slide-data-gen")
async def slide_data_gen(
    student_level: str = Form(...),
    document_url: str = Form(...),
    num_slides: int = Form(10),
    selected_topics: List[str] = Form([]),
    settings: Settings = Depends(get_settings)
):
    """Generate slide data for presentation"""
    try:
        start_time = time.time()
        print(f"\n🔍 Starting slide-data-gen for {student_level} level with {num_slides} slides")
        print(f"  Document URL: {document_url}")
        
        # Load prerequisites from file
        prerequisites_path = Path("data/metadata/prerequisites_dict.json")
        if not prerequisites_path.exists():
            raise HTTPException(
                status_code=404, 
                detail="Prerequisites data not found. Please run analyze/url or analyze/pdf first."
            )
            
        try:
            prerequisites_dict = load_json(prerequisites_path)
            print(f"✅ Loaded prerequisites with {len(prerequisites_dict)} topics")
        except Exception as e:
            print(f"❌ Error loading prerequisites: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to load prerequisites data: {str(e)}")
        
        # Filter prerequisites based on selected topics
        if selected_topics:
            filtered_prerequisites = {topic: prerequisites_dict[topic] 
                                    for topic in selected_topics 
                                    if topic in prerequisites_dict}
            print(f"🔍 Filtered prerequisites down to {len(filtered_prerequisites)} selected topics")
        else:
            filtered_prerequisites = prerequisites_dict
            print("⚠️ No topics selected, using all prerequisite topics")
        
        # Initialize Mistral client
        client = Mistral(api_key=settings.MISTRAL_API_KEY)

        slide_prompt_1 = f"""Analyze this research paper and create a structured outline for a PowerPoint presentation tailored for a {student_level} audience, where student_level is either 'PhD researcher', 'Master's student', or 'Undergrad student'. Follow these guidelines:

Pictures should only contain figure number, example (figure 2)

Important: Once a formula is shown do not use it again

CONSIDER Pictures AND FORMULA AS IMAGES

DO NOT KEEP IMAGE/Pictures AND FORMULA IN SAME SLIDE

Total slides should be : {num_slides} **Always follow this, must not be more than this**

Title Slide: Create a concise, engaging title that captures the paper's essence.

This is the first slide and a catchy subtitle. Bullets must be empty in this

Agenda slide: The second slide which contains all subtitles of other slides as its bullet points

Then State the paper's main research question and significance

This acts like the title

Highlight key background information relevant to the {student_level} audience

For each section of the paper:

Create 'Section Overview' slide with:

Section title: bullet points summarizing key concepts in content

Any critical formulas, using LaTeX notation

Followed with in depth explanation slides that only if necessary and also break each point in section overview into one slide explaining in depth:

Explain complex ideas using analogies or visualizations

Break down important formulas step-by-step

Highlight connections to prerequisite knowledge for the {student_level}

Do this for all the chosen topics

Make sure to

Present key findings with supporting data or graphs

Interpret results at a level appropriate for the {student_level}

With great amount of details not missing any key information

Discussion & Implications (1-2 slides):

Outline the paper's main conclusions

Discuss potential applications or future research directions

Relate findings to broader field context for the {student_level}

Key Takeaways slide:

List 3-5 main points to remember, tailored to the {student_level}'s level

Further Reading slide:

Suggest 3-4 related papers or resources appropriate for the {student_level}

For each slide, provide:

A clear, concise headline as the subtitle

Bullet points for main content (5-7 per slide)

Notes on any visuals, charts, or diagrams to include from the paper mention the figure number

Adjust the depth and complexity of the content based on the {student_level}, ensuring the presentation is informative and engaging for the specified audience level.

Make sure to include the mathematical formulas where ever necessary

Give the output in a json format and a dictionary tagging formuala and its name in json
"""
        
        # Create slide prompt
        slide_prompt_2 = f"""Analyze this research paper and create a structured outline for a PowerPoint presentation tailored for a {student_level} audience, where student_level is either 'PhD researcher', 'Master's student', or 'Undergrad student'. Follow these guidelines:

Pictures should only contain figure number, example (figure 2)

ONLY INCLUDE THE TOPICS IN {filtered_prerequisites} FOR THE SLIDES

Important: Once a formula is shown do not use it again

CONSIDER Pictures AND FORMULA AS IMAGES

DO NOT KEEP IMAGE/Pictures AND FORMULA IN SAME SLIDE

Total slides should be : {num_slides} **Always follow this, must not be more than this**

Title Slide: Create a concise, engaging title that captures the paper's essence.

This is the first slide and a catchy subtitle. Bullets must be empty in this

Agenda slide: The second slide which contains all subtitles of other slides as its bullet points

Then State the paper's main research question and significance

This acts like the title

Highlight key background information relevant to the {student_level} audience

For each section of the paper:

Create 'Section Overview' slide with:

Section title: bullet points summarizing key concepts in content

Any critical formulas, using LaTeX notation

Followed with in depth explanation slides that only if necessary and also break each point in section overview into one slide explaining in depth:

Explain complex ideas using analogies or visualizations

Break down important formulas step-by-step

Highlight connections to prerequisite knowledge for the {student_level}

Do this for all the chosen topics

Make sure to

Present key findings with supporting data or graphs

Interpret results at a level appropriate for the {student_level}

With great amount of details not missing any key information

Discussion & Implications (1-2 slides):

Outline the paper's main conclusions

Discuss potential applications or future research directions

Relate findings to broader field context for the {student_level}

Key Takeaways slide:

List 3-5 main points to remember, tailored to the {student_level}'s level

Further Reading slide:

Suggest 3-4 related papers or resources appropriate for the {student_level}

For each slide, provide:

A clear, concise headline as the subtitle

Bullet points for main content (5-7 per slide)

Notes on any visuals, charts, or diagrams to include from the paper mention the figure number

Adjust the depth and complexity of the content based on the {student_level}, ensuring the presentation is informative and engaging for the specified audience level.

Make sure to include the mathematical formulas where ever necessary

Give the output in a json format and a dictionary tagging formuala and its name in json
"""        
        # Define the messages for the chat
        if len(filtered_prerequisites) == 0:
            slide_prompt = slide_prompt_1
        else:
            slide_prompt = slide_prompt_2

        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": slide_prompt
                    },
                    {
                        "type": "document_url",
                        "document_url": document_url
                    }
                ]
            }
        ]

        api_start = time.time()
        print("📞 Calling Mistral API for slide generation...")
        try:
            chat_response = client.chat.parse(
                model=settings.MODEL_NAME,
                messages=messages,
                response_format=Ppt
            )
            api_time = time.time() - api_start
            print(f"✅ API call completed in {api_time:.2f}s")
            
            slides = chat_response.choices[0].message.content
        except Exception as e:
            print(f"❌ Error in Mistral API call: {e}")
            raise HTTPException(status_code=500, detail=f"Error generating slides content: {str(e)}")
        
        input_slides_path = Path(config.TEMPLATE_METADATA_DIR) / "slides_data.json"
        try:
            save_json(slides, input_slides_path)
            print(f"✅ Saved slides data to {input_slides_path}")
        except Exception as e:
            print(f"❌ Error saving slides data: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to save slides data: {str(e)}")
        
        total_time = time.time() - start_time
        print(f"✅ Slide data generation completed in {total_time:.2f}s")
        
        return {"message": "Slide data generated successfully", "path": str(input_slides_path)}
    
    except HTTPException as he:
        # Re-raise HTTP exceptions
        raise he
    except Exception as e:
        print(f"🚨 Unexpected error in slide-data-gen: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error generating slide data: {str(e)}")

@app.post("/process-slides-data")
async def process_slides_data():
    """Process slides data to include formulas and figures"""
    try:
        start_time = time.time()
        print("\n🔥 Starting /process-slides-data")

        input_slides_path = Path(config.TEMPLATE_METADATA_DIR) / "slides_data.json"
        output_slides_path = Path(config.TEMPLATE_METADATA_DIR) / "updated_slides_data.json"
        figures_metadata_path = Path(config.FIGURES_DIR) / "figures_metadata.json"

        if not input_slides_path.exists():
            raise HTTPException(status_code=404, detail=f"Missing slides_data.json. Make sure slide-data-gen was called successfully.")

        # Print content of metadata directory for debugging
        print(f"📂 Contents of metadata directory:")
        metadata_dir = Path(config.TEMPLATE_METADATA_DIR)
        for file_path in metadata_dir.glob("*"):
            print(f"  - {file_path.name} ({file_path.stat().st_size} bytes)")

        try:
            slides_data = load_json(input_slides_path)
        except Exception as e:
            print(f"❌ Error loading slides_data.json: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to load slides_data.json: {str(e)}")
        
        # Check if figures metadata exists
        image_data = {}
        if figures_metadata_path.exists():
            try:
                image_data = load_json(figures_metadata_path)
                print(f"✅ Loaded figures metadata with {len(image_data)} entries")
            except Exception as e:
                print(f"⚠️ Warning: Could not load figures metadata: {e}")
        else:
            print(f"⚠️ Warning: Figures metadata file not found at {figures_metadata_path}")

        # Process formulas in parallel for better performance
        slides_data = process_formulas_parallel(slides_data)

        # Process pictures
        for slide in slides_data.get("content", []):
            if "picture" in slide and slide["picture"]:
                for i in range(len(slide["picture"])):
                    key = slide["picture"][i].capitalize()
                    if key in image_data:
                        slide["picture"][i] = image_data[key]
                    else:
                        print(f"⚠️ Warning: Missing image data for key '{key}'")
                        slide["picture"][i] = ""

        # Save the updated JSON
        try:
            save_json(slides_data, output_slides_path)
            print(f"✅ Saved updated slides data to {output_slides_path}")
        except Exception as e:
            print(f"❌ Error saving updated slides data: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to save updated slides data: {str(e)}")

        elapsed = time.time() - start_time
        print(f"✅ Slides processed in {elapsed:.2f}s")
        return {"message": "Slides processed and updated JSON saved.", "path": str(output_slides_path)}

    except HTTPException as he:
        print(f"🚨 HTTP Error {he.status_code}: {he.detail}")
        raise he
    except Exception as e:
        print(f"🚨 Unexpected Error: {type(e).__name__}: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

@app.post("/enhace-slides-agent")
async def enhance_slides_agent(
    settings: Settings = Depends(get_settings)
):
    print(f"🔍 Starting enhancer agent parsing")
    client = Mistral(api_key=settings.MISTRAL_API_KEY)
    
    execution_agent_id = settings.ENHANCE_AGENT_ID
    json_dir = Path("data/metadata")
    slides_data_path = json_dir / "updated_slides_data.json"
    slides_data = load_json(slides_data_path)
    print(f"📊 Loaded slides_data with {len(slides_data.get('content', []))} slides")
    if not settings.ENHANCE_AGENT_ID:
            print("⚠️ EXECUTION_AGENT_ID not found, falling back to standard chat completion")
    query=f"""
Understand the current slides data as provided:
{slides_data}

And more data to the slides and maintain the same json format as the output

For slides with formulas, explain them in technical terms rather than giving examples of usage
"""
    def run_analysis_execution_agent(query):
        """
        Sends a user query to a Python agent and returns the response.

        Args:
            query (str): The user query to be sent to the Python agent.

        Returns:
            str: The response content from the Python agent.
        """
        try:
            response = client.agents.complete(
                agent_id= execution_agent_id,
                messages = [
                    {
                        "role": "user",
                        "content":  query
                    },
                ],
                # response_format=Ppt
            )
            result = response.choices[0].message.content
            return result
        except Exception as e:
            print(f"Request failed: {e}. Please check your request.")
            return None
        
    
    enhance_agent = run_analysis_execution_agent(query)
    
    # Add validation to prevent saving null data
    try:
        data = extract_json(enhance_agent)
        
        # Validate the data is not null/empty and has expected structure
        if data is None:
            print("⚠️ Warning: Agent returned null data. Not updating slides.")
            return {"message": "Agent returned null data. Original slides kept unchanged."}
            
        # Check if data has expected structure (content array)
        if not isinstance(data, dict) or 'content' not in data or not isinstance(data['content'], list) or len(data['content']) == 0:
            print("⚠️ Warning: Agent returned invalid data format. Not updating slides.")
            print(f"Data structure received: {type(data)}")
            return {"message": "Agent returned invalid data format. Original slides kept unchanged."}
            
        # Verify content has some expected fields
        sample_slide = data['content'][0]
        required_fields = ['title', 'subtitle']
        missing_fields = [field for field in required_fields if field not in sample_slide]
        
        if missing_fields:
            print(f"⚠️ Warning: Agent output is missing required fields: {missing_fields}. Not updating slides.")
            return {"message": f"Agent output is missing required fields: {missing_fields}. Original slides kept unchanged."}
        
        # If we got here, data seems valid, so save it
        print(f"✅ Agent returned valid data with {len(data['content'])} slides. Saving to {slides_data_path}")
        save_json(data, slides_data_path)
        print(f"✅ Enhancer agent parsing completed successfully")
        return {"message": "Slides data enhanced and saved successfully"}
    except Exception as e:
        print(f"❌ Error processing agent output: {str(e)}")
        traceback.print_exc()
        return {"message": f"Error enhancing slides: {str(e)}. Original slides kept unchanged."}
    """Enhance slides data using execution agent"""

@app.post("/execution-agent-parsing")
async def execution_agent_parsing(
    template_name: str = Form("template.pptx"),
    settings: Settings = Depends(get_settings)
):
    """Parse slides data using execution agent or model to create presentation structure"""
    try:
        start_time = time.time()
        print(f"🔍 Starting execution agent parsing with template: {template_name}")
        
        # Check for template with case-insensitive matching
        upload_dir = Path("data/upload")
        template_path = None
        
        # Try exact match first
        exact_match = upload_dir / template_name
        if exact_match.exists():
            template_path = exact_match
            print(f"📌 Using exact template match: {template_name}")
        else:
            # Try case-insensitive matching
            for file in upload_dir.glob("*.pptx"):
                if file.name.lower() == template_name.lower():
                    template_path = file
                    print(f"📌 Found template with case-insensitive match: {file.name}")
                    break
        
        if not template_path:
            raise HTTPException(
                status_code=400, 
                detail=f"Template file {template_name} not found. Available templates: {', '.join([f.name for f in upload_dir.glob('*.pptx')])}"
            )
            
        template_dir = template_path
        json_dir = Path("data/metadata")
        slides_data_path = json_dir / "updated_slides_data.json"
        layout_data_path = json_dir / "processed_layout.json"
        output_path = json_dir / "execution_agent.json"

        # Check for required files
        if not slides_data_path.exists():
            raise HTTPException(status_code=404, detail="updated_slides_data.json not found")
        if not layout_data_path.exists():
            raise HTTPException(status_code=404, detail="processed_layout.json not found")

        presentation = Presentation(template_dir)
        distinct_layout = [layout.name for layout in presentation.slide_master.slide_layouts]
        print(f"📊 Available layouts in template: {distinct_layout}")

        try:
            slides_data = load_json(slides_data_path)
            
        except Exception as e:
            print(f"❌ Error loading slides_data.json: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to load slides_data.json: {str(e)}")
            
        try:
            slides_layout = load_json(layout_data_path)
            print(f"📊 Loaded layout data with {len(slides_layout)} layouts")
            
            # Debug: print the first layout structure
            if slides_layout and len(slides_layout) > 0:
                print(f"📋 Sample layout structure: {slides_layout[0]}")
        except Exception as e:
            print(f"❌ Error loading processed_layout.json: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to load processed_layout.json: {str(e)}")

        query=f"""
        Understand the current slides data as provided:
        {slides_data}

        Now create a new json with layouts chosen from {distinct_layout} for each of the slide. The json contains the placeholders.name_placeholders.index from {slides_layout} as the key and the content as the values.

        Using the fields in the slides data namely title, subtitle, text, formula_images and picture to the placeholder of the layout selected for the slide.

        Give new titles for each slide and remove the subtitle

        Treat both formulas and pictures as images when assigning the layout

        If the content has bullet points, then be creative in choosing layouts which contain texts and title type

        Always add the layout name in the json chosen from the {distinct_layout}

        Add \ n in the placeholders to create new lines

        ## Do not add your own custom placeholders in the slide layout, use only the ones provided in the layout_details.json

        ##final output must in JSON only"""
        
        client = Mistral(api_key=settings.MISTRAL_API_KEY)
        
        # Check if EXECUTION_AGENT_ID is available
        if not settings.EXECUTION_AGENT_ID:
            print("⚠️ EXECUTION_AGENT_ID not found, falling back to standard chat completion")
            # Fallback to standard chat completion
            response = client.chat.complete(
                model=settings.MODEL_NAME,
                messages=[{"role": "user", "content": query}],
                # response_format=PresentationData
            )
            agent_result = response.choices[0].message.content
        else:
            try:
                print(f"🤖 Using execution agent ID: {settings.EXECUTION_AGENT_ID}")
                response = client.agents.complete(
                    agent_id=settings.EXECUTION_AGENT_ID,
                    messages=[{"role": "user", "content": query}],
                    # response_format=PresentationData
                )
                agent_result = response.choices[0].message.content
            except Exception as agent_error:
                print(f"⚠️ Agent execution failed, falling back to standard chat: {agent_error}")
                # Fallback to standard chat completion
                response = client.chat.complete(
                    model=settings.MODEL_NAME,
                    messages=[{"role": "user", "content": query}],
                    # response_format=PresentationData
                )
                agent_result = response.choices[0].message.content
        
        print("Got the agent Response")
        

        execution_agent_json = extract_json(agent_result)
        
        if not execution_agent_json:
            raise HTTPException(status_code=400, detail="Failed to extract valid JSON from agent response")
            
        # Ensure the execution_agent_json has the required structure
        if "slides" not in execution_agent_json:
            print("⚠️ Adding missing 'slides' key to execution_agent_json")
            execution_agent_json = {"slides": execution_agent_json}
            
        if len(execution_agent_json["slides"]) == 0:
            raise HTTPException(status_code=400, detail="Generated JSON contains no slides")
            
        # Validate the structure
        for slide in execution_agent_json["slides"]:
            if "slide_name" not in slide:
                raise HTTPException(status_code=400, detail=f"Missing slide_name in slide: {slide}")
            if slide["slide_name"] not in distinct_layout:
                print(f"⚠️ Warning: Slide layout '{slide['slide_name']}' not found in template layouts")
            if "placeholders" not in slide:
                raise HTTPException(status_code=400, detail=f"Missing placeholders in slide: {slide}")

        save_json(execution_agent_json, output_path)
        
        elapsed = time.time() - start_time
        print(f"✅ Execution agent parsing completed in {elapsed:.2f}s")

        return {"message": "Execution agent parsing completed", "path": str(output_path)}
    
    except HTTPException as he:
        # Re-raise HTTP exceptions
        raise he
    except Exception as e:
        print(f"🚨 Error in execution agent parsing: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error in execution agent parsing: {str(e)}")

@app.post("/generate-presentation")
async def generate_presentation_from_execution_json(
    template_name: str = Form("template.pptx"),
    execution_json_filename: str = Form("execution_agent.json"),
    output_ppt_filename: str = Form("modified_presentation.pptx"),
    processed_layout_filename: str = Form("processed_layout.json"),
    optimize_images: bool = Form(True),
    remove_slides_count: int = Form(13)
):
    """Generate the final presentation based on execution agent JSON"""
    print(f"📌 Starting generate-presentation with template: {template_name}")
    try:
        start_time = time.time()
        
        # Prepare directories
        output_dir = Path("data/output")
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Find the template file with case-insensitive matching
        upload_dir = Path("data/upload")
        template_path = None
        
        # Try exact match first
        exact_match = upload_dir / template_name
        if exact_match.exists():
            template_path = exact_match
            print(f"📌 Using exact template match: {template_name}")
        else:
            # Try case-insensitive matching
            for file in upload_dir.glob("*.pptx"):
                if file.name.lower() == template_name.lower():
                    template_path = file
                    print(f"📌 Found template with case-insensitive match: {file.name}")
                    break
        
        if not template_path:
            raise HTTPException(
                status_code=400, 
                detail=f"Template file {template_name} not found. Available templates: {', '.join([f.name for f in upload_dir.glob('*.pptx')])}"
            )
            
        # Set up paths
        template_dir = template_path
        json_dir = Path("data/metadata")
        execution_agent_json = json_dir / execution_json_filename
        output_ppt_path = output_dir / output_ppt_filename

        # Check if files exist
        if not execution_agent_json.exists():
            raise HTTPException(status_code=404, detail=f"{execution_json_filename} not found. Make sure execution-agent-parsing was called.")
            
        # Print content of metadata directory for debugging
        print(f"📂 Contents of metadata directory:")
        for file_path in json_dir.glob("*"):
            print(f"  - {file_path.name} ({file_path.stat().st_size} bytes)")
            
        # Load the presentation
        try:
            prs = Presentation(template_dir)
            print(f"✅ Presentation loaded from {template_dir}")
        except Exception as e:
            print(f"❌ Error loading template: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to load presentation template: {str(e)}")
            
        # Get available layouts for debugging
        layout_names = [layout.name for layout in prs.slide_master.slide_layouts]
        print(f"📋 Available layouts in template: {layout_names}")
        
        # Load the JSON data
        try:
            json_data = load_json(execution_agent_json)
            print(f"✅ Loaded execution agent JSON")
        except Exception as e:
            print(f"❌ Error loading execution agent JSON: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to load execution agent JSON: {str(e)}")
            
        # Validate the JSON structure
        if "slides" not in json_data:
            print("⚠️ Missing 'slides' key in execution_agent.json, checking if it's a direct array...")
            if isinstance(json_data, list):
                print("✅ Found array, wrapping in slides object")
                json_data = {"slides": json_data}
            else:
                print(f"❌ Invalid JSON structure: {json_data}")
                raise HTTPException(status_code=400, detail="Invalid JSON structure in execution_agent.json")
                
        # Debug print the JSON structure
        print(f"📊 JSON contains {len(json_data.get('slides', []))} slides")
        for i, slide in enumerate(json_data.get("slides", [])):
            print(f"  - Slide {i+1}: {slide.get('slide_name', 'UNKNOWN')} with {len(slide.get('placeholders', {}))} placeholders")

        def get_layout_by_name(prs_obj, layout_name):
            for layout in prs_obj.slide_layouts:
                if layout.name == layout_name:
                    return layout
            return None

        # Function to optimize images if needed
        def optimize_image(image_path, max_width=1280, quality=85):
            if not optimize_images:
                return image_path
                
            try:
                image = Image.open(image_path)
                # Only resize if image is larger than max dimensions
                if image.width > max_width:
                    ratio = max_width / image.width
                    new_height = int(image.height * ratio)
                    image = image.resize((max_width, new_height), Image.LANCZOS)
                
                # Save optimized version
                optimized_path = f"{os.path.splitext(image_path)[0]}_optimized.jpg"
                image.save(optimized_path, "JPEG", quality=quality, optimize=True)
                print(f"🔍 Optimized image: {os.path.basename(image_path)}")
                return optimized_path
            except Exception as e:
                print(f"⚠️ Image optimization failed: {e}")
                return image_path

        # Add slides based on JSON data
        slides_added = 0
        print(f"🔄 Creating slides from JSON data...")
        
        for slide_index, slide_data in enumerate(json_data.get("slides", [])):
            layout_name = slide_data.get("layout")
            if not layout_name:
                try:
                    layout_name = slide_data.get("slide_name")
                except Exception as e:
                    print(f"⚠️ Missing slide_name in slide {slide_index+1}, skipping")
                    continue
                
            layout = get_layout_by_name(prs, layout_name)

            if layout:
                try:
                    slide = prs.slides.add_slide(layout)
                    slides_added += 1
                    print(f"✅ Added slide {slide_index+1} with layout '{layout_name}'")
                    
                    if "placeholders" not in slide_data:
                        print(f"⚠️ No placeholders found in slide {slide_index+1}")
                        continue
                        
                    for placeholder_name, content in slide_data["placeholders"].items():
                        name_parts = placeholder_name.split("_")
                        if len(name_parts) < 2:
                            print(f"⚠️ Invalid placeholder name format: {placeholder_name}, expected name_index")
                            continue
                            
                        name = name_parts[0]
                        index = name_parts[-1]
                        
                        try:
                            # Find placeholder by index
                            placeholder_found = False
                            idx = int(index)
                            for shape in slide.shapes:
                                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == idx:
                                    placeholder_found = True
                                    
                                    # Handle text content
                                    if shape.has_text_frame:
                                        if content is None:
                                            shape.text_frame.text = ""
                                        elif isinstance(content, list):
                                            shape.text_frame.clear()
                                            for item in content:
                                                paragraph = shape.text_frame.add_paragraph()
                                                paragraph.text = str(item)
                                        else:
                                            # Check if this is a formula path
                                            if isinstance(content, str) and ("formulas" in content or "Formula" in content) and content.endswith(".png"):
                                                try:
                                                    # This is a formula image, insert it
                                                    formula_path = content
                                                    # Ensure full path
                                                    if not os.path.exists(formula_path) and not formula_path.startswith("data/"):
                                                        formula_path = f"data/formulas/{os.path.basename(formula_path)}"
                                                    
                                                    if os.path.exists(formula_path):
                                                        print(f"🔤 Inserting formula image at {formula_path}")
                                                        shape.insert_picture(formula_path)
                                                    else:
                                                        print(f"⚠️ Formula image not found: {formula_path}")
                                                        # Try to re-render if possible
                                                        shape.text_frame.text = "Formula image missing"
                                                except Exception as e:
                                                    print(f"⚠️ Error inserting formula image: {e}")
                                                    shape.text_frame.text = ""  # Don't show the path
                                            else:
                                                shape.text_frame.text = str(content)
                                    
                                    # Handle picture content
                                    if name.startswith("Picture") or "Picture" in name or shape.placeholder_format.type == 18:  # 18 is picture type
                                        try:
                                            if content and isinstance(content, str):
                                                # Path normalization - replace backslashes with forward slashes
                                                pic_path = content.replace("\\", "/")
                                                
                                                # Check if this is actually an image path
                                                if not pic_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                                                    shape.text_frame.text = str(content)
                                                    continue
                                                    
                                                # Make sure it has the full path
                                                if not os.path.exists(pic_path):
                                                    if not pic_path.startswith("data/"):
                                                        if "formula" in pic_path.lower():
                                                            pic_path = f"data/formulas/{os.path.basename(pic_path)}"
                                                        elif "figure" in pic_path.lower() or "img" in pic_path.lower():
                                                            pic_path = f"data/figures/{os.path.basename(pic_path)}"
                                                
                                                # Final existence check with detailed error
                                                if os.path.exists(pic_path):
                                                    try:
                                                        print(f"🖼️ Inserting picture: {pic_path}")
                                                        # Direct insertion attempt
                                                        shape.insert_picture(pic_path)
                                                    except Exception as img_error:
                                                        print(f"❌ Error inserting image: {str(img_error)}")
                                                        # Try to insert in text frame as fallback
                                                        if hasattr(shape, 'text_frame'):
                                                            shape.text_frame.text = f"[Image: {os.path.basename(pic_path)}]"
                                                else:
                                                    print(f"⚠️ Image file not found: {pic_path} (exists check failed)")
                                                    if hasattr(shape, 'text_frame'):
                                                        shape.text_frame.text = f"[Missing image: {os.path.basename(pic_path)}]"
                                        except Exception as e:
                                            print(f"⚠️ Error handling picture placeholder: {e}")
                                            traceback.print_exc()
                                            if hasattr(shape, 'text_frame'):
                                                shape.text_frame.text = "[Image error]"
                            
                            if not placeholder_found:
                                print(f"⚠️ Placeholder with index {index} not found in slide {slide_index+1}")
                        
                        except ValueError as e:
                            print(f"⚠️ Error processing placeholder {placeholder_name}: {e}")
                            continue
                
                except Exception as e:
                    print(f"❌ Error adding slide {slide_index+1}: {e}")
                    traceback.print_exc()
                    continue
            else:
                print(f"⚠️ Layout '{layout_name}' not found in the template. Available layouts: {layout_names}")

        if slides_added == 0:
            raise HTTPException(status_code=400, detail="No slides could be added to the presentation. Check log for details.")
            
        print(f"📊 Total slides added: {slides_added}")
            
        # Save temporary presentation
        temp_ppt_path = "temp.pptx"
        prs.save(temp_ppt_path)
        print(f"💾 Saved temporary presentation to {temp_ppt_path}")

        prs_1 = Presentation(template_dir)
        total__template_slides = len(prs_1.slides)

        # Load the presentation again to process slides
        prs = Presentation(temp_ppt_path)
        total_slides = len(prs.slides)
        
        # Only remove slides if there are enough slides
        remove_slides_count = total__template_slides 
        if total_slides > remove_slides_count:
            slides_to_remove = remove_slides_count
            print(f"🧹 Removing {slides_to_remove} slides from the beginning")
            
            for _ in range(slides_to_remove):
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            print(f"⚠️ Not removing slides: only {total_slides} slides available, need more than {remove_slides_count}")

        # Save the final presentation
        prs.save(str(output_ppt_path))
        print(f"✅ Final presentation saved to {output_ppt_path}")
        
        # Clean up temp file
        os.remove(temp_ppt_path)
        
        elapsed = time.time() - start_time
        print(f"✅ Presentation generated in {elapsed:.2f}s")

        return {"message": "Presentation generated successfully", "path": str(output_ppt_path)}
    
    except Exception as e:
        print(f"❌ Error generating presentation: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error generating presentation: {str(e)}")

@app.get("/download-presentation")
async def download_presentation(filename: str = "modified_presentation.pptx"):
    """Download the generated presentation"""
    file_path = Path("data/output") / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Presentation file not found")
    return FileResponse(file_path, filename=filename)

@app.post("/cleanup-data")
async def cleanup_data(should_clean: bool = Form(False)):
    """
    Optional cleanup endpoint - only to be called manually when needed, 
    not automatically after successful generation
    """
    try:
        if should_clean:
            # Clean directories only if explicitly requested
            clean_directories(["data/figures", "data/formulas"])
            return {"message": "Successfully cleaned up data directories"}
        else:
            return {"message": "No cleanup performed - set should_clean=true to perform cleanup"}
    except Exception as e:
        print(f"❌ Error cleaning up data: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error cleaning up data: {str(e)}")

# Podcast related constants
PODCAST_SYSTEM_PROMPT = """
You are a world-class podcast producer tasked with transforming the provided input text into an engaging and informative podcast script. The input may be unstructured or messy, sourced from PDFs or web pages. Your goal is to extract the most interesting and insightful content for a compelling podcast discussion.

# Steps to Follow:

1. **Analyze the Input:**
   Carefully examine the text, identifying key topics, points, and interesting facts or anecdotes that could drive an engaging podcast conversation. Disregard irrelevant information or formatting issues.

2. **Brainstorm Ideas:**
   In the `<scratchpad>`, creatively brainstorm ways to present the key points engagingly. Consider:
   - Analogies, storytelling techniques, or hypothetical scenarios to make content relatable
   - Ways to make complex topics accessible to a general audience
   - Thought-provoking questions to explore during the podcast
   - Creative approaches to fill any gaps in the information

3. **Craft the Dialogue:**
   Develop a natural, conversational flow between the host (Jane) and the guest speaker (the author or an expert on the topic). Incorporate:
   - The best ideas from your brainstorming session
   - Clear explanations of complex topics
   - An engaging and lively tone to captivate listeners
   - A balance of information and entertainment

   Rules for the dialogue:
   - The host (Jane) always initiates the conversation and interviews the guest
   - Include thoughtful questions from the host to guide the discussion
   - Incorporate natural speech patterns, including occasional verbal fillers (e.g., "um," "well," "you know")
   - Allow for natural interruptions and back-and-forth between host and guest
   - Ensure the guest's responses are substantiated by the input text, avoiding unsupported claims
   - Maintain a PG-rated conversation appropriate for all audiences
   - Avoid any marketing or self-promotional content from the guest
   - The host concludes the conversation

4. **Summarize Key Insights:**
   Naturally weave a summary of key points into the closing part of the dialogue. This should feel like a casual conversation rather than a formal recap, reinforcing the main takeaways before signing off.

5. **Maintain Authenticity:**
   Throughout the script, strive for authenticity in the conversation. Include:
   - Moments of genuine curiosity or surprise from the host
   - Instances where the guest might briefly struggle to articulate a complex idea
   - Light-hearted moments or humor when appropriate
   - Brief personal anecdotes or examples that relate to the topic (within the bounds of the input text)

6. **Consider Pacing and Structure:**
   Ensure the dialogue has a natural ebb and flow:
   - Start with a strong hook to grab the listener's attention
   - Gradually build complexity as the conversation progresses
   - Include brief "breather" moments for listeners to absorb complex information
   - End on a high note, perhaps with a thought-provoking question or a call-to-action for listeners

IMPORTANT RULE: Each line of dialogue should be no more than 100 characters (e.g., can finish within 5-8 seconds)

Give the output in a json format
"""

QUESTION_MODIFIER = "PLEASE ANSWER THE FOLLOWING QN:"
TONE_MODIFIER = "TONE: The tone of the podcast should be"
LANGUAGE_MODIFIER = "OUTPUT LANGUAGE <IMPORTANT>: The the podcast should be"
LENGTH_MODIFIERS = {
    "Short (1-2 min)": "Keep the podcast brief, around 1-2 minutes long.",
    "Medium (3-5 min)": "Aim for a moderate length, about 3-5 minutes.",
    "Long (10-15 min)": "Aim for a longer podcast, around 10-15 minutes.",
}

# Podcast schema classes
class DialogueItem(BaseModel):
    speaker: str  # "Host (Jane)" or "Guest"
    text: str

class PodcastDialogue(BaseModel):
    scratchpad: str
    name_of_guest: str
    dialogue: List[DialogueItem]

@app.post("/generate-podcast/")
async def api_generate_podcast(
    request: dict,
    settings: Settings = Depends(get_settings)
):
    """Generate a podcast from the document URL"""
    try:
        print(f"🎙️ Starting podcast generation with request: {request}")
        
        document_url = request.get("document_url")
        prompt_modifiers = request.get("prompt_modifiers", {})
        
        if not document_url:
            raise HTTPException(status_code=400, detail="document_url is required")
            
        # Call the generate_podcast function from the imported module
        result = generate_podcast(document_url, prompt_modifiers)
        
        return result
        
    except Exception as e:
        print(f"❌ Error generating podcast: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error generating podcast: {str(e)}")

# Keep the old endpoint for compatibility
@app.post("/generate-podcast")
async def legacy_generate_podcast(
    question: Optional[str] = Form(None),
    tone: str = Form("Fun"),
    length: str = Form("Medium (3-5 min)"),
    language: str = Form("English"),
    settings: Settings = Depends(get_settings),
    document_url: Optional[str] = Form(None)
):
    """Legacy endpoint for generating a podcast from a document URL using form data"""
    print(f"🎙️ Starting legacy podcast generation")
    
    try:
        if not document_url:
            raise HTTPException(status_code=400, detail="document_url is required")
            
        # Convert form parameters to the new format
        prompt_modifiers = {}
        if question:
            prompt_modifiers["question"] = question
        if tone:
            prompt_modifiers["tone"] = tone
        if length:
            prompt_modifiers["length"] = length
        if language:
            prompt_modifiers["language"] = language
            
        # Call the generate_podcast function from the imported module
        result = generate_podcast(document_url, prompt_modifiers)
        
        return result
        
    except Exception as e:
        print(f"❌ Error generating podcast: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error generating podcast: {str(e)}")

@app.get("/podcast/{filename}")
async def get_podcast_file(filename: str):
    """Get podcast file by filename"""
    podcast_path = Path("podcast") / filename
    if not podcast_path.exists():
        raise HTTPException(status_code=404, detail=f"Podcast file '{filename}' not found")
    
    print(f"Serving podcast file: {podcast_path}")
    # Explicitly set media_type to audio/mpeg to ensure proper handling by browsers
    return FileResponse(
        path=str(podcast_path),
        media_type="audio/mpeg",
        filename=filename
    )

# Main application entry point
if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000) 